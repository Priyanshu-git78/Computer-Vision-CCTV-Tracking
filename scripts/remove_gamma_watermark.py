# Ultralytics 🚀 AGPL-3.0 License - https://ultralytics.com/license

import os
import sys
import pptx
from pptx import Presentation
import logging

# Ensure project root is in path
project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if project_root not in sys.path:
    sys.path.insert(0, project_root)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

PPTX_PATHS = [
    "/mnt/hdd/ultralytics/AI-Powered-Intelligent-Security-Ecosystem.pptx",
    "/mnt/hdd/ultralytics/project/AI-Powered-Intelligent-Security-Ecosystem.pptx"
]

def remove_watermark():
    for pptx_path in PPTX_PATHS:
        if not os.path.exists(pptx_path):
            logger.warning(f"File not found: {pptx_path}")
            continue

        logger.info(f"Removing Gamma watermark from: {pptx_path}")
        prs = Presentation(pptx_path)
        
        # Access Slide Layout 1 (index 1) which all slides use
        layout = prs.slide_layouts[1]
        
        # Find the shape with name="Image 0" or located at (8.78, 5.30)
        shape_to_remove = None
        for shape in layout.shapes:
            left = shape.left / 914400
            top = shape.top / 914400
            w = shape.width / 914400
            h = shape.height / 914400
            
            # Match by position or name
            if (shape.name == "Image 0" and abs(left - 8.78) < 0.1 and abs(top - 5.30) < 0.1):
                shape_to_remove = shape
                break
                
        if shape_to_remove is not None:
            layout.shapes._spTree.remove(shape_to_remove._element)
            logger.info("Watermark shape successfully removed from layout.")
            prs.save(pptx_path)
            logger.info(f"Saved: {pptx_path}")
        else:
            logger.warning("Gamma watermark shape not found in layout.")

if __name__ == "__main__":
    remove_watermark()
