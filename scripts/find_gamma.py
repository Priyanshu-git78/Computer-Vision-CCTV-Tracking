# Ultralytics 🚀 AGPL-3.0 License - https://ultralytics.com/license

import os
from pptx import Presentation

PPTX_PATH = "/mnt/hdd/ultralytics/AI-Powered-Intelligent-Security-Ecosystem.pptx"

def find_gamma():
    if not os.path.exists(PPTX_PATH):
        print("File not found!")
        return

    prs = Presentation(PPTX_PATH)
    
    # 1. Check Slide Masters
    for idx, master in enumerate(prs.slide_masters):
        for s_idx, shape in enumerate(master.shapes):
            if shape.has_text_frame and "gamma" in shape.text_frame.text.lower():
                print(f"[Master {idx}] Shape {s_idx} text contains 'gamma': {shape.text_frame.text}")
            if "gamma" in shape.name.lower():
                print(f"[Master {idx}] Shape {s_idx} name contains 'gamma': {shape.name}")
                
    # 2. Check Slide Layouts
    for idx, layout in enumerate(prs.slide_layouts):
        for s_idx, shape in enumerate(layout.shapes):
            if shape.has_text_frame and "gamma" in shape.text_frame.text.lower():
                print(f"[Layout {idx}] Shape {s_idx} text contains 'gamma': {shape.text_frame.text}")
            if "gamma" in shape.name.lower():
                print(f"[Layout {idx}] Shape {s_idx} name contains 'gamma': {shape.name}")

    # 3. Check Slides
    for idx, slide in enumerate(prs.slides):
        for s_idx, shape in enumerate(slide.shapes):
            if shape.name and "gamma" in shape.name.lower():
                print(f"[Slide {idx+1}] Shape {s_idx} name: {shape.name}")
            if shape.click_action and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                addr = shape.click_action.hyperlink.address.lower()
                if "gamma" in addr:
                    print(f"[Slide {idx+1}] Shape {s_idx} hyperlink: {shape.click_action.hyperlink.address}")

if __name__ == "__main__":
    find_gamma()
