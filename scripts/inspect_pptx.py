# Ultralytics 🚀 AGPL-3.0 License - https://ultralytics.com/license

import os
import pptx
from pptx import Presentation

pptx_path = "/mnt/hdd/ultralytics/AI-Powered-Intelligent-Security-Ecosystem.pptx"

def inspect():
    if not os.path.exists(pptx_path):
        print("File not found!")
        return

    prs = Presentation(pptx_path)
    print(f"Total slides: {len(prs.slides)}")
    
    for idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {idx+1} ---")
        
        # Check shapes for text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        print(f"[Text] {text}")
            
            # Check for images
            if shape.shape_type == 13: # Picture shape
                image = shape.image
                print(f"[Image] format={image.ext}, size={len(image.blob)} bytes")

if __name__ == "__main__":
    inspect()
