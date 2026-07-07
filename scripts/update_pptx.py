# Ultralytics 🚀 AGPL-3.0 License - https://ultralytics.com/license

import os
import sys
import docx
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import logging

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

PPTX_PATH = "/mnt/hdd/ultralytics/AI-Powered-Intelligent-Security-Ecosystem.pptx"
IMAGE_PATH = "/mnt/hdd/ultralytics/project/outputs/train/tuned_run/confusion_matrix.png"

def update_presentation():
    if not os.path.exists(PPTX_PATH):
        logger.error(f"Presentation not found: {PPTX_PATH}")
        return

    logger.info(f"Opening presentation: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)

    # 1. Update Slide 9 (Index 8) - Experimental Results
    slide9 = prs.slides[8]
    logger.info("Updating Slide 9 contents...")

    # Clear existing shapes on Slide 9
    for shape in list(slide9.shapes):
        slide9.shapes._spTree.remove(shape._element)

    # Add Title
    title_box = slide9.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Experimental Evaluation & Results"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(26, 54, 93) # Deep Blue

    # Add Subtitle
    sub_box = slide9.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.4))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Empirical measurements from model training, tracking comparison, and system orchestration."
    p_sub.font.size = Pt(13)
    p_sub.font.italic = True
    p_sub.font.color.rgb = RGBColor(113, 128, 150) # Slate Gray

    # Add Column 1: Custom Model Training
    col1_box = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2.8), Inches(3.6))
    tf_col1 = col1_box.text_frame
    tf_col1.word_wrap = True
    
    p1 = tf_col1.paragraphs[0]
    p1.text = "Custom Model Training"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(43, 108, 176)
    
    bullets1 = [
        "• Person Detector (YOLOv8n)",
        "  - mAP50: 0.3733 on COCO-Person",
        "  - Optimized for lightweight tracking",
        "• Tuned Fire Detector (YOLOv8n)",
        "  - mAP50: 0.6383 | mAP50-95: 0.4611",
        "  - Fine-tuned with optimized augmentation",
        "• Specialist Fire (YOLOv8s)",
        "  - mAP50: 0.6950 | mAP50-95: 0.3953",
        "  - Higher capacity validation model"
    ]
    for b in bullets1:
        p = tf_col1.add_paragraph()
        p.text = b
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(45, 55, 72)

    # Add Column 2: Tracking & System Efficiency
    col2_box = slide9.shapes.add_textbox(Inches(3.4), Inches(1.5), Inches(2.8), Inches(3.6))
    tf_col2 = col2_box.text_frame
    tf_col2.word_wrap = True
    
    p2 = tf_col2.paragraphs[0]
    p2.text = "System Performance"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(43, 108, 176)
    
    bullets2 = [
        "• Tracking (TownCentre Benchmark)",
        "  - ByteTrack: 100.0% count accuracy",
        "  - BoT-SORT: 93.8% count accuracy",
        "• Standard vs. Adaptive System",
        "  - Frame Rate: 30.0 -> 12.6 FPS (-58%)",
        "  - VRAM Saving: 1.8GB -> 0.9GB (-50%)",
        "  - Event Recall: 100.0% (No missed events)",
        "  - Achieved substantial compute reduction",
        "    without sacrificing security recall"
    ]
    for b in bullets2:
        p = tf_col2.add_paragraph()
        p.text = b
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(45, 55, 72)

    # Add Column 3: Image (Confusion Matrix)
    if os.path.exists(IMAGE_PATH):
        logger.info(f"Adding confusion matrix image to Slide 9: {IMAGE_PATH}")
        slide9.shapes.add_picture(IMAGE_PATH, Inches(6.3), Inches(1.6), width=Inches(3.2), height=Inches(3.2))
    else:
        logger.warning(f"Image not found for Slide 9: {IMAGE_PATH}")

    # 2. Update Slide 10 (Index 9) - Replace placeholder text
    slide10 = prs.slides[9]
    logger.info("Checking Slide 10 for progress placeholder...")
    for shape in slide10.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if "training experiments in progress" in paragraph.text.lower() or "results pending" in paragraph.text.lower():
                    logger.info("Replacing Slide 10 placeholder text.")
                    paragraph.text = "Training and evaluation successfully completed. Fine-tuned model weights and results have been validated and uploaded to Hugging Face."
                    paragraph.font.size = Pt(14)
                    paragraph.font.italic = True
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(34, 139, 34) # Green to show success

    # Save presentation
    prs.save(PPTX_PATH)
    logger.info(f"Successfully saved updated presentation to: {PPTX_PATH}")

    # Copy to project folder as well
    project_pptx = "/mnt/hdd/ultralytics/project/AI-Powered-Intelligent-Security-Ecosystem.pptx"
    import shutil
    shutil.copy2(PPTX_PATH, project_pptx)
    logger.info(f"Copied updated presentation to: {project_pptx}")

if __name__ == "__main__":
    update_presentation()
